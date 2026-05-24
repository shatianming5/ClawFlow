from __future__ import annotations

import textwrap
from pathlib import Path


SLIDES = [
    ("ClawFlow", "面向下一代个人智能体的轻量级 Agent Runtime", "docs/assets/diagrams/architecture.png"),
    ("从 Chatbot 到 AgentOS", "Agent 需要可执行、可记忆、可恢复、可治理的基础设施。", "docs/assets/diagrams/agentos_kernel.png"),
    ("现有问题", "Demo 化严重、缺少状态恢复、权限不可控、执行不可观测、插件和记忆难管理。", None),
    ("项目目标", "能执行、能记忆、能恢复、能治理、能扩展、能支撑真实应用构建。", None),
    ("总体架构", "Gateway + AgentOS Kernel + Tool Sandbox + Memory + Trace + Plugin + Applications。", "docs/assets/diagrams/architecture.png"),
    ("AgentOS Kernel", "Runtime Kernel、State Manager、Event Bus、Scheduler、Policy Engine。", "docs/assets/diagrams/agentos_kernel.png"),
    ("Agent Runtime", "Planner 生成 Plan，Executor 执行 Step，Tool Registry 提供能力边界。", "docs/assets/diagrams/workflow.png"),
    ("Workflow + Checkpoint", "Step Graph、Retry 边界、Resume、Replay、Workflow Export。", "docs/assets/diagrams/workflow.png"),
    ("Tool Sandbox", "BaseTool 元数据、risk_level、schema、shell 白名单。", "docs/assets/diagrams/tool_governance.png"),
    ("Permission Governance", "allow / ask / deny、Human-in-the-loop、Audit Log。", "docs/assets/screenshots/web_governance.png"),
    ("Human Approval Queue", "Web/API/CLI approve or deny pending high-risk tool calls with trace and audit updates.", "docs/assets/screenshots/web_approvals.png"),
    ("Memory Layer", "SQLite 长期记忆、keyword search、hit count、vector interface。", "docs/assets/screenshots/web_memory.png"),
    ("Observability", "Trace Timeline、Metrics、Error、Cost estimation、Replay。", "docs/assets/screenshots/web_trace_timeline.png"),
    ("Prompt Registry + Cost Dashboard", "Prompt templates, usage count, estimated tokens, cost and failure analysis.", "docs/assets/screenshots/web_cost_dashboard.png"),
    ("Tool Usage Heatmap", "Aggregate real tool_call trace events to show runtime behavior and infrastructure observability.", "docs/assets/screenshots/web_tool_usage_heatmap.png"),
    ("Multi-channel Gateway", "CLI、FastAPI、Web UI 共享同一个 Runtime。", "docs/assets/screenshots/web_home.png"),
    ("Plugin System", "Manifest-driven Plugin Registry and MCP-like connector expansion.", "docs/assets/diagrams/plugin_system.png"),
    ("Multi-agent Collaboration", "ManagerAgent、ResearchAgent、ToolAgent、CriticAgent、ReportAgent。", "docs/assets/diagrams/multi_agent_topology.png"),
    ("RAG Module", "Document ingestion、chunk、retrieval trace、grounded answer。", "docs/assets/diagrams/rag_pipeline.png"),
    ("From Demo to Example Application", "Example Applications 不是孤立脚本，而是 Runtime 验证工作负载。", "docs/assets/diagrams/example_application_stack.png"),
    ("Application 1: Research Assistant", "生成项目摘要、报告大纲、TODO，并写入 trace/memory/checkpoint。", "docs/assets/screenshots/cli_research_application.png"),
    ("Application 2: Personal Assistant", "生成 daily_plan.md 并写入长期记忆。", "docs/assets/screenshots/web_memory.png"),
    ("Application 3: Safe Tool Call", "高风险删除请求转为 dry-run，记录 audit。", "docs/assets/screenshots/web_audit_log.png"),
    ("Application 4: Multi-agent / RAG", "多角色协作与文档检索回答均通过统一 Runtime。", "docs/assets/screenshots/web_multi_agent.png"),
    ("Benchmark", "真实 Runtime 任务生成 latency、success rate、tool calls、trace events。", "docs/assets/figures/benchmark_latency.png"),
    ("开源协议与社区规范", "Apache-2.0、CONTRIBUTING、SECURITY、GOVERNANCE、ROADMAP。", None),
    ("对比与差异化", "从 demo.py 到可复用开发框架；从黑盒执行到可观测 Runtime。", None),
    ("开发者框架", "SDK wrapper、Application Template Generator、Tool Template Generator，支持二次开发。", "docs/assets/screenshots/web_applications.png"),
    ("未来展望", "MCP、云端部署、移动端入口、企业知识库、多模型路由、插件市场。", "docs/assets/screenshots/web_applications.png"),
    ("总结", "从 Prompt Demo 到 Agent Runtime；从孤立样例到统一 Runtime 驱动的基础设施。", "docs/assets/screenshots/web_benchmark.png"),
]


def write_outline_and_notes() -> None:
    outline = ["# ClawFlow PPT Outline", ""]
    notes = ["# Slide Notes", ""]
    for idx, (title, body, image) in enumerate(SLIDES, 1):
        outline.append(f"{idx}. {title} - {body}")
        notes.append(f"## Slide {idx}: {title}\n\n{body}\n\n演讲提示：强调 ClawFlow 不是 ChatGPT API 套壳，而是 Agent Runtime / AgentOS Infrastructure。")
        if image:
            outline.append(f"   - Asset: `{image}`")
    Path("slides").mkdir(exist_ok=True)
    Path("slides/ppt_outline.md").write_text("\n".join(outline) + "\n", encoding="utf-8")
    Path("slides/slide_notes.md").write_text("\n\n".join(notes) + "\n", encoding="utf-8")


def generate_pptx() -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:
        Path("slides/ClawFlow_presentation_failed.md").write_text(f"PPTX generation failed: {exc}\n", encoding="utf-8")
        return
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for idx, (title, body, image) in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(14, 17, 22)
        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(6.8), Inches(0.65))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(30 if idx != 1 else 42)
        p.font.bold = True
        p.font.color.rgb = RGBColor(238, 242, 246)
        body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(5.5), Inches(4.8))
        tfb = body_box.text_frame
        tfb.word_wrap = True
        for line in textwrap.wrap(body, width=38):
            para = tfb.add_paragraph() if tfb.text else tfb.paragraphs[0]
            para.text = line
            para.font.size = Pt(21)
            para.font.color.rgb = RGBColor(203, 213, 225)
            para.space_after = Pt(8)
        tag = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(5.5), Inches(0.35))
        tag.text_frame.text = "Agent Runtime · AgentOS Kernel · Workflow Orchestration · Permission Governance · Trace Replay"
        tag.text_frame.paragraphs[0].font.size = Pt(10)
        tag.text_frame.paragraphs[0].font.color.rgb = RGBColor(73, 210, 180)
        if image and Path(image).exists():
            try:
                slide.shapes.add_picture(image, Inches(6.25), Inches(1.05), width=Inches(6.55), height=Inches(5.45))
            except Exception:
                pass
        else:
            panel = slide.shapes.add_shape(1, Inches(6.35), Inches(1.25), Inches(5.9), Inches(4.7))
            panel.fill.solid()
            panel.fill.fore_color.rgb = RGBColor(23, 27, 34)
            panel.line.color.rgb = RGBColor(43, 50, 61)
            text = slide.shapes.add_textbox(Inches(6.7), Inches(2.0), Inches(5.2), Inches(3.5))
            text.text_frame.text = "ClawFlow\nUnified Runtime\nReal State\nReplayable Trace"
            for para in text.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                para.font.size = Pt(25)
                para.font.color.rgb = RGBColor(238, 242, 246)
    prs.save("slides/ClawFlow_presentation.pptx")


def generate_pdf() -> None:
    try:
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfgen import canvas
    except Exception as exc:
        Path("slides/ClawFlow_presentation_pdf_failed.md").write_text(f"PDF generation failed: {exc}\n", encoding="utf-8")
        return
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    c = canvas.Canvas("slides/ClawFlow_presentation.pdf", pagesize=landscape(A4))
    width, height = landscape(A4)
    for idx, (title, body, image) in enumerate(SLIDES, 1):
        c.setFillColorRGB(0.055, 0.067, 0.086)
        c.rect(0, 0, width, height, fill=1, stroke=0)
        c.setFillColorRGB(0.93, 0.95, 0.97)
        c.setFont("STSong-Light", 28)
        c.drawString(42, height - 58, title)
        c.setFont("STSong-Light", 15)
        y = height - 105
        for line in textwrap.wrap(body, width=42):
            c.drawString(48, y, line)
            y -= 24
        c.setFillColorRGB(0.29, 0.82, 0.71)
        c.setFont("STSong-Light", 10)
        c.drawString(48, 30, f"Slide {idx} · ClawFlow Agent Runtime")
        c.showPage()
    c.save()


def main() -> None:
    write_outline_and_notes()
    generate_pptx()
    generate_pdf()
    print("generated slides/ClawFlow_presentation.*")


if __name__ == "__main__":
    main()
